from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
IMG_DIR = ROOT / "output" / "supplied-screenshots" / "新建文件夹"
TEST_DIR = ROOT / "国产测试工具"
OUT_DIR = ROOT / "output" / "ppt-assets"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = ROOT / "output" / "AI智慧课程平台-项目展示成品稿.pptx"


W, H = 13.333, 7.5

COLORS = {
    "bg": "F3F9FF",
    "paper": "FFFFFF",
    "panel": "F7FBFF",
    "ink": "0D2238",
    "muted": "526B86",
    "line": "C8E3FA",
    "grid": "E1F0FC",
    "teal": "0B8EA8",
    "teal2": "E0F7FB",
    "blue": "1D6FE8",
    "blue2": "E7F1FF",
    "cyan": "13B7D8",
    "cyan2": "E5FAFF",
    "gold": "D89A25",
    "gold2": "FFF4DB",
    "red": "D95D6A",
    "red2": "FFECEF",
    "green": "20A779",
    "green2": "E8FAF3",
}


def rgb(hex_color):
    return RGBColor.from_string(hex_color)


def image_path(name):
    return IMG_DIR / f"屏幕截图 2026-07-03 {name}.png"


IMAGES = {
    "map": image_path("193851"),
    "map_modal": image_path("193906"),
    "ai_chat": image_path("193934"),
    "chapter": image_path("193944"),
    "course_card": image_path("193952"),
    "task": image_path("194000"),
    "diagnosis": image_path("194020"),
    "battle": image_path("194049"),
    "db_table": image_path("194310"),
    "submission": image_path("194436"),
    "battle_ai": image_path("194530"),
    "modal": image_path("194630"),
    "teacher_home": image_path("194701"),
    "course_list": image_path("194719"),
    "knowledge_graph": image_path("194745"),
    "course_detail": image_path("195251"),
    "ability_map": image_path("195537"),
    "score": image_path("200059"),
    "analytics": image_path("200117"),
    "table": image_path("200810"),
    "student_profile": image_path("200913"),
    "questions": image_path("200925"),
    "exam": image_path("200933"),
    "bank": image_path("200939"),
    "resource": image_path("201711"),
    "test_tool": TEST_DIR / "屏幕截图 2026-06-26 185717.png",
    "test_report": TEST_DIR / "屏幕截图 2026-06-26 194141.png",
}


def crop_cover(src, ratio, out_name):
    out = OUT_DIR / out_name
    im = Image.open(src).convert("RGB")
    w, h = im.size
    # Keep the full screenshot visible. Earlier versions used cover-cropping to
    # fill fixed frames, but that removed important UI information.
    target_w = max(1600, w)
    target_h = int(target_w / ratio)
    if target_h < h:
        target_h = h
        target_w = int(target_h * ratio)
    canvas = Image.new("RGB", (target_w, target_h), "#" + COLORS["bg"])
    scale = min(target_w / w, target_h / h)
    new_size = (int(w * scale), int(h * scale))
    resized = im.resize(new_size, Image.Resampling.LANCZOS)
    canvas.paste(resized, ((target_w - new_size[0]) // 2, (target_h - new_size[1]) // 2))
    canvas.save(out, quality=94)
    return out


def set_font(run, size=16, bold=False, color="ink"):
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color] if color in COLORS else color)


def add_text(slide, text, x, y, w, h, size=16, bold=False, color="ink",
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color)
    return box


def add_multiline(slide, lines, x, y, w, h, size=14, color="muted"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        set_font(run, size, False, color)
    return box


def add_shape(slide, kind, x, y, w, h, fill, line=None, radius=False, transparency=0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else kind
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COLORS[fill] if fill in COLORS else fill)
    shp.fill.transparency = transparency
    if line:
        shp.line.color.rgb = rgb(COLORS[line] if line in COLORS else line)
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def add_card(slide, x, y, w, h, fill="panel", line="line"):
    return add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h, fill, line, radius=True)


def add_pill(slide, text, x, y, w, h, fill, color="ink", size=12):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h, fill, None, radius=True)
    add_text(slide, text, x, y + 0.05, w, h - 0.08, size=size, bold=True, color=color,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_img(slide, src, x, y, w, h, radius=False, border=True):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x - 0.03, y - 0.03, w + 0.06, h + 0.06, "paper", "line")
    pic = slide.shapes.add_picture(str(src), Inches(x), Inches(y), Inches(w), Inches(h))
    if border:
        pic.line.color.rgb = rgb(COLORS["line"])
        pic.line.width = Pt(0.55)
    return pic


def add_title(slide, section, title, subtitle=None):
    add_pill(slide, section, 0.55, 0.35, 1.3, 0.34, "cyan2", "teal", 10)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.55, 0.73, 0.08, 0.72, "cyan", None)
    add_text(slide, title, 0.55, 0.8, 7.3, 0.52, size=26, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.57, 1.32, 8.3, 0.35, size=11.5, color="muted")


def add_footer(slide, idx):
    add_text(slide, f"{idx:02d}", 12.37, 7.05, 0.42, 0.2, size=9, color="muted",
             align=PP_ALIGN.RIGHT)


def slide_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["bg"])
    # Subtle technology-style grid and circuit traces.
    for x in [0.7, 2.4, 4.1, 5.8, 7.5, 9.2, 10.9, 12.6]:
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(0), Inches(0.006), Inches(H))
        line.fill.solid()
        line.fill.fore_color.rgb = rgb(COLORS["grid"])
        line.fill.transparency = 45
        line.line.fill.background()
    for y in [0.65, 1.85, 3.05, 4.25, 5.45, 6.65]:
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(y), Inches(W), Inches(0.006))
        line.fill.solid()
        line.fill.fore_color.rgb = rgb(COLORS["grid"])
        line.fill.transparency = 50
        line.line.fill.background()
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.32, 0.22, 1.15, 0.045, "cyan", None, transparency=12)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 11.65, 7.05, 1.05, 0.045, "blue", None, transparency=20)
    for x, y in [(1.55, 0.24), (11.55, 7.07), (12.8, 0.9)]:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, y, 0.07, 0.07, "cyan", None, transparency=5)


def add_metric(slide, value, label, x, y, color="blue"):
    add_text(slide, value, x, y, 1.5, 0.42, size=25, bold=True, color=color)
    add_text(slide, label, x, y + 0.47, 1.5, 0.28, size=10.5, color="muted")


def add_labeled_note(slide, label, title, body, x, y, w, h, color="teal"):
    add_card(slide, x, y, w, h)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, 0.12, h, f"{color}2", None)
    add_pill(slide, label, x + 0.24, y + 0.18, 0.9, 0.28, f"{color}2", color, 8.5)
    add_text(slide, title, x + 0.24, y + 0.56, w - 0.5, 0.22, size=12.8, bold=True)
    add_text(slide, body, x + 0.24, y + 0.86, w - 0.5, h - 1.0, size=10.5, color="muted")


def add_relation_strip(slide, items, x, y, w, h, color="teal"):
    gap = 0.12
    item_w = (w - gap * (len(items) - 1)) / len(items)
    for i, (label, detail) in enumerate(items):
        xx = x + i * (item_w + gap)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, xx, y, item_w, h, f"{color}2", color, radius=True)
        add_text(slide, label, xx + 0.12, y + 0.16, item_w - 0.24, 0.22, size=11.5, bold=True,
                 color=color, align=PP_ALIGN.CENTER)
        add_text(slide, detail, xx + 0.12, y + 0.48, item_w - 0.24, 0.18, size=8.8,
                 color="muted", align=PP_ALIGN.CENTER)


def add_data_loop(slide, items, x, y, w, h):
    cx = x + w / 2
    cy = y + h / 2
    positions = [
        (x + 0.1, y + 0.18),
        (x + w - 1.75, y + 0.18),
        (x + w - 1.75, y + h - 0.8),
        (x + 0.1, y + h - 0.8),
    ]
    for (label, detail, color), (xx, yy) in zip(items, positions):
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, xx, yy, 1.65, 0.62, f"{color}2", color, radius=True)
        add_text(slide, label, xx + 0.1, yy + 0.12, 1.45, 0.18, size=10.8, bold=True,
                 color=color, align=PP_ALIGN.CENTER)
        add_text(slide, detail, xx + 0.1, yy + 0.36, 1.45, 0.14, size=7.9, color="muted",
                 align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, cx - 0.72, cy - 0.43, 1.44, 0.86, "paper", "line")
    add_text(slide, "数据闭环", cx - 0.55, cy - 0.12, 1.1, 0.18, size=12.0, bold=True,
             color="teal", align=PP_ALIGN.CENTER)

def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # Pre-cropped assets
    crops = {
        "map_wide": crop_cover(IMAGES["map"], 1.78, "map_wide.jpg"),
        "diagnosis_wide": crop_cover(IMAGES["diagnosis"], 1.78, "diagnosis_wide.jpg"),
        "battle_wide": crop_cover(IMAGES["battle"], 2.05, "battle_wide.jpg"),
        "battle_ai_wide": crop_cover(IMAGES["battle_ai"], 1.78, "battle_ai_wide.jpg"),
        "teacher_home": crop_cover(IMAGES["teacher_home"], 1.75, "teacher_home.jpg"),
        "knowledge": crop_cover(IMAGES["knowledge_graph"], 1.75, "knowledge.jpg"),
        "ability": crop_cover(IMAGES["ability_map"], 1.75, "ability.jpg"),
        "analytics": crop_cover(IMAGES["analytics"], 1.75, "analytics.jpg"),
        "profile": crop_cover(IMAGES["table"], 1.75, "profile.jpg"),
        "test": crop_cover(IMAGES["test_tool"], 1.55, "test.jpg"),
        "test_report": crop_cover(IMAGES["test_report"], 1.55, "test_report.jpg"),
    }

    # 1 Cover
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_shape(s, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 4.6, H, "cyan2", None, transparency=10)
    add_shape(s, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 4.6, 0, 0.03, H, "cyan", None, transparency=12)
    add_text(s, "AI 智慧课程平台", 0.72, 0.98, 4.0, 0.72, size=32, bold=True)
    add_text(s, "知识图谱驱动的爬塔学习系统", 0.74, 1.86, 3.55, 0.4, size=17, bold=True, color="teal")
    add_text(s, "AI Agent + KingbaseES + 国产化测试工具", 0.74, 2.3, 3.55, 0.35, size=14, color="muted")
    add_multiline(s, [
        "面向 Python 程序设计课程，平台把知识点、题目、作答记录、能力画像和教学干预组织成一条可追踪的数据链。",
        "核心特色是“爬塔学习”：学生沿知识图谱路线完成诊断、战斗、复盘和通关，学习过程具有连续目标和即时反馈。"
    ], 0.76, 3.0, 3.45, 1.6, size=12.8, color="ink")
    add_text(s, "项目展示汇报", 7.4, 1.35, 3.1, 0.42, size=28, bold=True, color="blue")
    add_text(s, "AI 智慧课程平台", 7.42, 1.95, 3.25, 0.34, size=17, bold=True)
    add_text(s, "爬塔学习系统 · AI 诊断推荐 · 国产数据库 · 国产化测试", 7.42, 2.42, 4.3, 0.28, size=12.5, color="muted")
    add_shape(s, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 7.42, 3.08, 3.35, 0.045, "cyan", None, transparency=8)
    add_relation_strip(s, [
        ("课程", "Python"),
        ("学习", "爬塔闯关"),
        ("智能", "AI Agent"),
        ("国产化", "KingbaseES"),
    ], 6.55, 4.55, 5.7, 0.78, "blue")
    add_text(s, "轻量科技风展示稿", 8.0, 5.65, 2.35, 0.24, size=12.5, color="teal", align=PP_ALIGN.CENTER)
    add_footer(s, 1)

    # 2 Overall design
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "总体设计", "爬塔学习系统的支撑结构",
              "将界面体验、AI 诊断、知识图谱、国产数据库和教师决策拆成相互支撑的业务层。")
    levels = [
        ("教学决策", "学情分析 / 干预建议", "blue2", "blue"),
        ("AI 诊断", "错因分析 / 推荐路线", "teal2", "teal"),
        ("房间挑战", "诊断房 / 战斗房 / 补给房", "gold2", "gold"),
        ("知识图谱", "前置关系 / 能力映射", "green2", "green"),
        ("国产数据底座", "KingbaseES V9", "red2", "red"),
    ]
    for i, (t, d, fill, c) in enumerate(levels):
        y = 0.88 + i * 1.05
        x = 5.55 + i * 0.23
        w_level = 5.8 - i * 0.46
        add_card(s, x, y, w_level, 0.78, fill=fill, line=c)
        add_text(s, t, x + 0.32, y + 0.18, 1.35, 0.2, size=13.5, bold=True, color=c)
        add_text(s, d, x + 1.82, y + 0.19, w_level - 2.2, 0.2, size=11.5, color="ink")
    add_relation_strip(s, [
        ("学生端", "闯关学习"),
        ("爬塔系统", "诊断与战斗"),
        ("AI 反馈", "提示与诊断"),
        ("教师端", "分析与干预"),
    ], 5.35, 6.32, 6.35, 0.74, "teal")
    add_card(s, 0.75, 1.95, 3.75, 3.75, fill="panel", line="line")
    add_text(s, "设计说明", 1.08, 2.28, 1.8, 0.28, size=18, bold=True, color="blue")
    add_multiline(s, [
        "学生端呈现为游戏化爬塔体验，核心数据仍由后端统一维护。",
        "AI Agent 参与推荐与诊断，但通关、节点、任务状态由后端校验。",
        "KingbaseES 保存课程、任务、答题、画像与分析数据，支撑国产化部署。",
        "教师端使用这些过程数据完成风险识别和教学干预。"
    ], 1.08, 2.9, 2.9, 1.75, 12.3)
    add_footer(s, 2)

    # 3 Positioning
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "项目定位", "解决“学什么、怎么练、谁需要干预”三个问题",
              "平台价值不在功能堆叠，而在把学生学习过程变成可计算、可推荐、可干预的数据闭环。")
    add_card(s, 0.65, 1.95, 3.25, 3.95)
    add_text(s, "痛点", 0.95, 2.25, 2.3, 0.35, size=18, bold=True)
    add_multiline(s, ["学生不知道下一步优先补哪个知识点。", "题目练习与知识掌握之间缺少可解释关联。", "教师只能看到结果分数，难以及时识别风险学生。", "AI 如果只做问答，无法真正进入教学流程。"], 0.95, 2.95, 2.55, 2.1, 12.4)
    add_card(s, 5.0, 1.55, 3.35, 4.75, fill="teal2", line="teal")
    add_text(s, "方案", 5.35, 1.95, 2.2, 0.35, size=19, bold=True, color="teal")
    add_multiline(s, ["知识图谱决定学习路线，突破单纯章节顺序的限制。", "诊断房判断是否已掌握，战斗房承担针对性训练。", "AI 根据错因、掌握度和路线状态参与推荐与解释。", "教师端把过程数据转化为风险提醒和干预建议。"], 5.35, 2.7, 2.55, 2.25, 12.4, "ink")
    add_card(s, 9.3, 1.95, 3.25, 3.95)
    add_text(s, "落地", 9.62, 2.25, 2.3, 0.35, size=18, bold=True)
    add_multiline(s, ["前端：学生端爬塔体验 + 教师端管理分析。", "后端：REST API 统一承载课程、任务、题库、画像等业务。", "Agentic：知识抽取、路线推荐、诊断反馈、教学建议。", "数据库：KingbaseES 支撑国产化部署与业务数据持久化。"], 9.62, 2.95, 2.55, 2.1, 12.4)
    add_footer(s, 3)

    # 3 Tower core
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "核心创新", "爬塔学习系统重构课程学习流程",
              "每个房间都有明确教学含义：诊断负责判断，战斗负责训练，补给负责修复，Boss 负责综合检验。")
    add_img(s, crops["map_wide"], 0.7, 1.85, 7.15, 4.65, radius=True)
    labels = [("塔", "Python 程序设计"), ("楼层", "知识点单元"), ("房间", "诊断 / 战斗 / 补给"), ("通关", "掌握度达标")]
    for i, (k, v) in enumerate(labels):
        x = 8.35 + (i % 2) * 2.05
        y = 1.95 + (i // 2) * 1.55
        add_card(s, x, y, 1.75, 1.15)
        add_text(s, k, x + 0.18, y + 0.18, 1.1, 0.24, size=13, bold=True, color="teal")
        add_text(s, v, x + 0.18, y + 0.58, 1.32, 0.28, size=12.2)
    add_multiline(s, ["设计重点：每个游戏元素都对应一个教学动作。", "节点状态来自后端持久化记录，避免前端只做静态展示。", "学生答题会同时影响战斗结算、知识点掌握度和后续路线。", "教师端看到的风险、薄弱点和建议，来自这些连续过程数据。"], 8.35, 4.95, 3.9, 1.55, 12.2)
    add_footer(s, 4)

    # 4 Knowledge route
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "路径生成", "个性化路线：AI 建议，后端校验，数据库留痕",
              "路线会根据学生掌握度、错题记录和节点状态持续调整，并保留完整过程记录。")
    add_img(s, crops["knowledge"], 0.75, 1.75, 6.45, 4.1, radius=True)
    add_card(s, 7.75, 1.7, 4.8, 4.2)
    add_text(s, "路线生成逻辑", 8.1, 2.02, 2.4, 0.32, size=18, bold=True)
    add_relation_strip(s, [
        ("知识关系", "前置依赖"),
        ("答题证据", "错题/耗时"),
        ("掌握状态", "节点进度"),
        ("能力画像", "成长趋势"),
    ], 8.1, 2.72, 4.0, 0.74, "blue")
    add_multiline(s, ["输入：知识点前置关系、能力点映射、历史错题、掌握度。", "AI：给出推荐路线和薄弱点修复顺序。", "后端：校验节点合法性、课程边界和通关状态。", "数据库：记录 run / node / attempt，保证路线可追溯。"], 8.1, 4.15, 3.95, 1.45, 12.2)
    add_footer(s, 5)

    # 5 Student map
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "学生端", "地图页承担“学习导航”和“状态解释”",
              "地图页集中呈现课程进度、知识点状态、房间类型和下一步行动。")
    add_img(s, crops["map_wide"], 0.65, 1.65, 8.1, 5.25, radius=True)
    add_card(s, 9.2, 1.78, 3.25, 4.95)
    add_text(s, "关键数据", 9.55, 2.12, 2.2, 0.3, 17, True)
    add_multiline(s, ["节点状态：未解锁、可进入、已通过、待复盘。", "房间类型：诊断、普通战斗、精英战斗、休息、商店。", "学习状态：等级、能量、金币、攻击、防御、经验。", "教学含义：让学生知道为什么进入这个节点、需要补什么。"], 9.55, 2.75, 2.45, 2.1, 12.2)
    add_metric(s, "等级 5", "学习成长状态", 9.55, 5.35, "teal")
    add_metric(s, "67%", "攻击能力", 11.0, 5.35, "gold")
    add_footer(s, 6)

    # 6 Diagnosis
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "学生端", "诊断房：用低成本测评决定学习分流",
              "诊断房承担路线分流作用：已掌握则快速通过，薄弱则进入针对性训练。")
    add_img(s, crops["diagnosis_wide"], 0.75, 1.65, 7.3, 4.75, radius=True)
    add_card(s, 8.55, 1.75, 3.95, 1.08, fill="green2", line="green")
    add_text(s, "通过策略", 8.85, 2.02, 1.1, 0.25, 15, True, "green")
    add_text(s, "诊断全对可直接通关", 10.0, 2.02, 1.65, 0.25, 15, True)
    add_card(s, 8.55, 3.15, 3.95, 1.08, fill="gold2", line="gold")
    add_text(s, "修复策略", 8.85, 3.42, 1.1, 0.25, 15, True, "gold")
    add_text(s, "薄弱点进入战斗房", 10.0, 3.42, 1.65, 0.25, 15, True)
    add_card(s, 8.55, 4.55, 3.95, 1.55)
    add_multiline(s, ["题目与知识点绑定，结果直接更新节点掌握状态。", "诊断结果影响路线、通关记录和后续训练安排。", "满分通过会记录 pass_method，薄弱诊断会触发修复节点。"], 8.85, 4.82, 2.95, 1.0, 12.3)
    add_footer(s, 7)

    # 7 Battle
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "学生端", "战斗房：把判题结果转化为可感知反馈",
              "战斗房让每一次作答都产生清晰后果：对、错、扣血、奖励、错题沉淀。")
    add_img(s, crops["battle_wide"], 0.65, 1.92, 8.35, 4.08, radius=True)
    add_card(s, 9.45, 1.85, 2.95, 4.75)
    add_text(s, "判题如何进入游戏", 9.75, 2.22, 2.0, 0.3, 16.5, True)
    add_multiline(s, ["题目难度决定敌人强度，知识点决定房间主题。", "学生提交答案后，后端返回判题结果和得分。", "正确答案转化为攻击收益，错误答案转化为生命值损失。", "错题、耗时和选项会沉淀到后续 AI 诊断与教师端分析。"], 9.75, 2.78, 2.25, 2.25, 11.8)
    add_pill(s, "即时反馈", 9.75, 5.58, 1.08, 0.35, "blue2", "blue", 10)
    add_pill(s, "错题沉淀", 10.98, 5.58, 1.08, 0.35, "red2", "red", 10)
    add_footer(s, 8)

    # 8 AI tutor
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "学生端", "AI 导师：只在关键学习节点介入",
              "AI 在学生卡住、答错或需要复盘时给出解释、提示和修复方向。")
    add_img(s, crops["battle_ai_wide"], 0.72, 1.7, 7.0, 4.55, radius=True)
    add_card(s, 8.15, 1.7, 4.25, 4.55)
    add_text(s, "AI 介入边界", 8.5, 2.03, 2.2, 0.28, 18, True)
    add_relation_strip(s, [
        ("题目解释", "理解题意"),
        ("知识提示", "定位概念"),
        ("错因诊断", "修复方向"),
    ], 8.5, 2.72, 3.35, 0.74, "teal")
    add_multiline(s, ["输入上下文包括题目、知识点、学生答案、历史薄弱点。", "输出包括提示、错因解释和复习建议，避免直接削弱练习价值。", "AI 状态可追踪：pending / success / failed / mock，便于区分真实 AI 诊断和降级结果。"], 8.5, 4.15, 3.35, 1.35, 12.1)
    add_footer(s, 9)

    # 9 Teacher overview
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "教师端", "首页聚合教学风险与待处理事项",
              "教师端首页把待处理任务、异常提醒和最近变化放在一起，帮助教师先处理最需要关注的问题。")
    add_img(s, crops["teacher_home"], 0.75, 1.62, 6.6, 4.25, radius=True)
    for i, (t, d, c) in enumerate([
        ("待处理队列", "把批改、逾期、临近截止统一收敛，减少教师查找成本", "blue"),
        ("异常提醒", "从任务提交和学情中提取薄弱知识点、低分与复核压力", "gold"),
        ("风险信号", "根据投入、表现、基础掌握和最近活动识别学生风险", "red"),
    ]):
        y = 1.8 + i * 1.35
        add_card(s, 8.0, y, 4.05, 0.95, fill=f"{c}2", line=c)
        add_text(s, t, 8.28, y + 0.18, 1.8, 0.25, 14.5, True, c)
        add_text(s, d, 9.85, y + 0.17, 1.95, 0.35, 10.6, False, "ink")
    add_footer(s, 10)

    # 10 Content org
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "教师端", "内容生产链：课程配置如何变成爬塔房间",
              "教师端配置的课程、题库和资源共同构成学生端爬塔系统的内容来源和测评依据。")
    add_img(s, IMAGES["course_list"], 0.62, 1.55, 5.9, 2.65, radius=True)
    add_img(s, IMAGES["resource"], 6.85, 1.55, 5.85, 2.65, radius=True)
    add_labeled_note(s, "内容源", "课程内容结构化", "课程、章节、资源先被组织成可维护的数据，再映射到学生端的学习节点。", 0.9, 4.55, 3.65, 1.15, "blue")
    add_labeled_note(s, "题目源", "题库服务房间挑战", "题目与知识点绑定后，可用于诊断房、普通战斗房、精英战斗房和 Boss 房。", 4.85, 4.55, 3.65, 1.15, "gold")
    add_labeled_note(s, "数据源", "数据反哺教学", "学生答题、通关、错题和任务提交数据会回流到教师端，用于分析与干预。", 8.8, 4.55, 3.65, 1.15, "teal")
    add_footer(s, 11)

    # 11 Ability map
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "教师端", "能力图谱：把课程目标拆成可计算指标",
              "能力点连接课程目标和知识点，学生画像中的成长趋势由这些指标持续累积。")
    add_img(s, crops["ability"], 0.68, 1.58, 7.3, 4.7, radius=True)
    add_card(s, 8.45, 1.9, 3.85, 3.95)
    add_metric(s, "787", "能力点", 8.85, 2.35, "blue")
    add_metric(s, "207", "知识点映射", 10.55, 2.35, "teal")
    add_metric(s, "100%", "覆盖率", 8.85, 3.55, "green")
    add_multiline(s, ["能力点用于描述学生能够完成的任务。", "知识点映射用于解释能力短板来自哪些具体内容。", "覆盖率用于衡量课程目标是否被题目和学习节点充分支撑。"], 8.85, 4.55, 2.8, 1.05, 12.0)
    add_footer(s, 12)

    # 12 Analytics
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "教师端", "学情分析：从单个错题上升到班级共性问题",
              "系统统计知识点错误、题型表现和高频错题，让教师知道下一节课该讲什么。")
    add_img(s, crops["analytics"], 0.62, 1.55, 7.4, 4.8, radius=True)
    add_card(s, 8.45, 1.75, 3.9, 4.4)
    add_text(s, "决策闭环", 8.8, 2.1, 1.6, 0.28, 18, True)
    add_data_loop(s, [
        ("作答数据", "逐题记录", "blue"),
        ("知识统计", "掌握度", "teal"),
        ("问题聚类", "共性错误", "gold"),
        ("教学建议", "干预动作", "green"),
    ], 8.65, 2.62, 3.25, 1.65)
    add_multiline(s, ["知识点层面：识别错误率高、掌握度低的内容。", "题目层面：定位高频错题和异常题目。", "题型层面：比较选择、填空、简答、编程等题型表现。", "教学层面：生成复习安排、专项练习和课堂讲解建议。"], 8.8, 4.2, 2.9, 1.25, 11.8)
    add_footer(s, 13)

    # 13 Student profile
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "教师端", "学生画像：把“分数”拆成可干预信号",
              "画像页并列呈现学习状态、投入、解题表现、基础掌握和成长趋势，便于教师定位干预对象。")
    add_img(s, crops["profile"], 0.65, 1.6, 8.05, 5.0, radius=True)
    add_card(s, 9.15, 1.92, 3.05, 4.3)
    add_multiline(s, ["学习状态：判断是否存在持续掉队风险。", "学习投入：结合活跃度和最近活动识别沉默学生。", "解题表现：区分不会做、粗心错和题型短板。", "基础掌握：直接指向薄弱知识点数量。", "教师动作：查看详情或分配个性化任务。"], 9.52, 2.28, 2.3, 2.35, 12.0)
    add_footer(s, 14)

    # 14 Database
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "国产数据库", "KingbaseES V9：为国产化部署做工程适配",
              "数据库设计围绕类型兼容、逻辑关联和业务扩展进行工程适配。")
    add_img(s, IMAGES["db_table"], 0.7, 1.72, 5.8, 3.95, radius=True)
    add_card(s, 7.05, 1.65, 5.35, 4.15)
    add_text(s, "工程设计要点", 7.4, 1.98, 2.4, 0.3, 18, True)
    add_multiline(s, ["采用 KingbaseES V9，使用 Spring profile 独立管理国产数据库连接。", "34 张核心业务表覆盖课程、任务、测评、画像、推荐和分析。", "业务 ID 以 VARCHAR 为主，规避国产数据库中整数/字符串比较差异。", "不依赖数据库外键强绑定，由服务层维护逻辑关联，方便后续扩展和迁移。"], 7.4, 2.58, 4.35, 1.75, 12.2)
    add_relation_strip(s, [
        ("课程", "course"),
        ("任务", "task"),
        ("答题", "answer"),
        ("画像", "profile"),
        ("分析", "analytics"),
    ], 7.4, 5.05, 4.2, 0.65, "teal")
    add_footer(s, 15)

    # 15 AI
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "AI 设计", "AI Agent 承担三类业务能力",
              "AI 在平台中承担内容结构化、学生诊断和教师决策辅助，不直接替代业务规则。")
    add_img(s, crops["ability"], 0.7, 1.78, 5.95, 3.85, radius=True)
    items = [
        ("内容侧", "从课程资料中抽取知识点，辅助生成能力点和知识图谱关系。"),
        ("学习侧", "结合当前题目和学生历史表现，提供提示、错因诊断和个性化推荐。"),
        ("教学侧", "对班级错误和薄弱点进行聚类，生成可执行的教学干预建议。"),
    ]
    for i, (t, d) in enumerate(items):
        y = 1.78 + i * 1.28
        add_card(s, 7.15, y, 4.9, 0.9)
        add_text(s, t, 7.48, y + 0.2, 0.9, 0.22, 14.5, True, "teal")
        add_text(s, d, 8.38, y + 0.21, 3.05, 0.22, 12.2, False)
    add_card(s, 7.15, 5.7, 4.9, 0.65, fill="teal2", line="teal")
    add_text(s, "设计原则：AI 给建议，后端做校验，数据库留痕，教师可追溯", 7.45, 5.91, 4.1, 0.2, 12.5, True, "teal")
    add_footer(s, 16)

    # 16 Testing + Summary
    s = prs.slides.add_slide(blank); slide_bg(s)
    add_title(s, "质量保障", "国产化测试工具 + 自动化测试覆盖关键链路",
              "测试重点放在高风险链路：接口稳定性、AI 服务降级、游戏资源、路由跳转和压力场景。")
    add_img(s, crops["test"], 0.72, 1.58, 5.4, 3.6, radius=True)
    add_img(s, crops["test_report"], 6.45, 1.58, 5.4, 3.6, radius=True)
    for i, (v, l, c) in enumerate([("319", "后端测试用例", "blue"), ("167", "前端测试用例", "teal"), ("0", "失败 / 错误", "green")]):
        add_metric(s, v, l, 1.1 + i * 2.0, 5.75, c)
    add_card(s, 7.25, 5.55, 4.35, 0.88, fill="paper")
    add_text(s, "WebRunner 用于录制、脚本、任务、负载和性能监控；单元测试覆盖后端 319 例、前端 167 例", 7.45, 5.83, 3.95, 0.28, 11.8, True)
    add_footer(s, 17)

    prs.save(OUT_FILE)
    print(OUT_FILE)


if __name__ == "__main__":
    build()
