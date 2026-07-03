from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_FILE = ROOT / "output" / "AI智慧课程平台-项目展示成品稿.pptx"
W, H = 13.333, 7.5

COLORS = {
    "bg": "F3F9FF",
    "paper": "FFFFFF",
    "panel": "F7FBFF",
    "ink": "0D2238",
    "muted": "526B86",
    "line": "C7E1F6",
    "grid": "E1F0FC",
    "blue": "1D6FE8",
    "blue2": "E7F1FF",
    "cyan": "13B7D8",
    "cyan2": "E5FAFF",
    "teal": "0B8EA8",
    "teal2": "DDF7F4",
    "gold": "D89A25",
    "gold2": "FFF4DB",
    "green": "20A779",
    "green2": "E8FAF3",
    "red": "D95D6A",
    "red2": "FFECEF",
}


def rgb(key):
    return RGBColor.from_string(COLORS.get(key, key))


def set_font(run, size=14, bold=False, color="ink"):
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def text(slide, value, x, y, w, h, size=14, bold=False, color="ink",
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    set_font(run, size, bold, color)
    return box


def multiline(slide, lines, x, y, w, h, size=12.2, color="muted", gap=5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = line
        set_font(run, size, False, color)
    return box


def shape(slide, kind, x, y, w, h, fill, line=None, radius=False, transparency=0):
    st = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else kind
    shp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.fill.transparency = transparency
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def card(slide, x, y, w, h, fill="panel", line="line"):
    return shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h, fill, line, radius=True)


def pill(slide, value, x, y, w, h, fill="cyan2", color="teal", size=10):
    shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h, fill, None, radius=True)
    text(slide, value, x, y + 0.045, w, h - 0.08, size=size, bold=True,
         color=color, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("bg")
    for x in [0.7, 2.4, 4.1, 5.8, 7.5, 9.2, 10.9, 12.6]:
        shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, 0, 0.006, H, "grid", None, transparency=45)
    for y in [0.65, 1.85, 3.05, 4.25, 5.45, 6.65]:
        shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, y, W, 0.006, "grid", None, transparency=50)
    shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.32, 0.22, 1.15, 0.045, "cyan", None, transparency=12)
    shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 11.65, 7.05, 1.05, 0.045, "blue", None, transparency=20)
    shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 1.55, 0.24, 0.07, 0.07, "cyan", None, transparency=5)
    shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 11.55, 7.07, 0.07, 0.07, "cyan", None, transparency=5)


def title(slide, section, heading, sub=None):
    pill(slide, section, 0.55, 0.34, 1.35, 0.34)
    shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.55, 0.72, 0.08, 0.72, "cyan", None)
    text(slide, heading, 0.72, 0.76, 8.2, 0.54, size=25.5, bold=True)
    if sub:
        text(slide, sub, 0.72, 1.32, 9.2, 0.34, size=11.5, color="muted")


def footer(slide, idx):
    text(slide, f"{idx:02d}", 12.35, 7.05, 0.42, 0.2, size=9, color="muted", align=PP_ALIGN.RIGHT)


def module_box(slide, label, body, x, y, w, h, color="blue"):
    card(slide, x, y, w, h, fill=f"{color}2" if f"{color}2" in COLORS else "panel", line=color)
    text(slide, label, x + 0.22, y + 0.18, w - 0.44, 0.24, size=13.2, bold=True, color=color)
    text(slide, body, x + 0.22, y + 0.56, w - 0.44, h - 0.7, size=10.5, color="ink")


def stack(slide, items, x, y, w, h):
    gap = 0.16
    ih = (h - gap * (len(items) - 1)) / len(items)
    for i, (label, detail, color) in enumerate(items):
        yy = y + i * (ih + gap)
        card(slide, x, yy, w, ih, fill=f"{color}2", line=color)
        text(slide, label, x + 0.28, yy + 0.18, 1.4, 0.22, size=13, bold=True, color=color)
        text(slide, detail, x + 1.78, yy + 0.18, w - 2.1, 0.22, size=11.3)


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    slides = []

    # 1 Cover
    s = prs.slides.add_slide(blank); bg(s)
    shape(s, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 4.7, H, "cyan2", None, transparency=8)
    shape(s, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 4.7, 0, 0.03, H, "cyan", None, transparency=10)
    text(s, "AI 智慧课程平台", 0.72, 1.15, 4.0, 0.62, size=31, bold=True)
    text(s, "基于知识图谱与 AI Agent 的爬塔式学习系统设计", 0.74, 2.02, 3.55, 0.72, size=16, bold=True, color="teal")
    multiline(s, ["游戏化学习", "个性化路径", "智能诊断", "国产数据库", "国产化测试"], 0.78, 3.25, 3.0, 1.2, 12.5, "ink")
    text(s, "项目设计汇报", 7.25, 1.75, 3.6, 0.46, size=30, bold=True, color="blue")
    text(s, "设计思路 · 模块划分 · 核心机制 · 推广应用", 7.28, 2.45, 4.1, 0.3, size=14.5, color="muted")
    shape(s, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 7.28, 3.08, 3.8, 0.05, "cyan", None)
    module_box(s, "展示边界", "操作演示与界面效果放在视频答辩环节；PPT 聚焦系统设计、模块逻辑和推广价值。", 7.25, 4.0, 4.4, 1.05, "blue")
    footer(s, 1)

    # 2 Background
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "问题定义", "项目背景与问题定义", "从学生学习困难、教师教学决策和系统智能化三个维度定义项目价值。")
    module_box(s, "学生端问题", "学生常常不知道下一步优先补哪个知识点，题目练习与知识掌握之间缺少可解释关联。", 0.8, 2.0, 3.7, 2.3, "blue")
    module_box(s, "教师端问题", "教师面对大量作答数据时，难以及时识别班级共性问题、个体风险和需要干预的薄弱点。", 4.85, 2.0, 3.7, 2.3, "teal")
    module_box(s, "系统端问题", "AI、知识图谱和数据库需要形成稳定闭环，才能让推荐、诊断和教学建议具备可追溯依据。", 8.9, 2.0, 3.7, 2.3, "green")
    card(s, 1.6, 5.1, 10.1, 0.9, fill="paper", line="cyan")
    text(s, "核心目标：把课程学习过程转化为可组织、可计算、可推荐、可干预的数据闭环。", 2.0, 5.42, 9.2, 0.24, size=15, bold=True, color="blue", align=PP_ALIGN.CENTER)
    footer(s, 2)

    # 3 Overall
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "总体设计", "知识组织 - 学习挑战 - 数据沉淀 - 智能反馈 - 教学干预", "平台围绕学习过程闭环展开设计，而不是只围绕课程管理功能堆叠。")
    stack(s, [
        ("知识组织", "课程、章节、知识点、能力点和前置关系", "blue"),
        ("学习挑战", "爬塔节点、诊断房、战斗房、补给房和 Boss 房", "teal"),
        ("数据沉淀", "作答、节点状态、掌握度、错题、能力变化和行为记录", "green"),
        ("智能反馈", "知识抽取、路线推荐、错因诊断和教学建议生成", "gold"),
        ("教学干预", "学情分析、风险提醒、学生画像和精准任务分配", "red"),
    ], 1.2, 1.95, 10.9, 4.5)
    footer(s, 3)

    # 4 Modules
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "模块划分", "五个业务模块与一个 AI Agent 支撑模块", "模块划分对应数据闭环中的不同职责，便于后续扩展和国产化部署。")
    modules = [
        ("课程内容与知识组织", "课程、章节、资源、知识点、知识关系、能力点"),
        ("学习任务与过程跟踪", "学习任务、节点状态、任务提交、答题过程、通关记录"),
        ("测评与成果评价", "题库、测验、作答结果、AI 批改、评价记录"),
        ("学生画像与个性化学习", "掌握度、能力分数、成长历史、推荐内容、成就系统"),
        ("学情分析与教学决策", "班级表现、共性错误、风险学生、教学报告、干预建议"),
        ("AI Agent 支撑模块", "知识图谱构建、智能问答、个性化推荐、主观评价、教学分析"),
    ]
    for i, (a, b) in enumerate(modules):
        x = 0.75 + (i % 3) * 4.15
        y = 1.9 + (i // 3) * 2.05
        module_box(s, a, b, x, y, 3.65, 1.55, ["blue", "teal", "green", "gold", "red", "blue"][i])
    footer(s, 4)

    # 5 Tower
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "核心创新", "爬塔式学习系统", "将 Python 课程抽象为学习塔，将知识点组织为楼层和节点，将题目练习转化为房间挑战。")
    stack(s, [
        ("塔", "对应一门课程或完整学习主题", "blue"),
        ("章节区域", "对应课程阶段，如基础语法、数据处理、综合实践", "teal"),
        ("楼层节点", "对应具体知识点或能力点", "green"),
        ("房间类型", "对应诊断、训练、复盘、补给和综合检验", "gold"),
        ("通关条件", "由掌握度、答题表现和节点状态共同决定", "red"),
    ], 0.95, 1.9, 5.4, 4.55)
    card(s, 7.05, 2.0, 4.95, 3.7, fill="paper", line="cyan")
    text(s, "设计目的", 7.45, 2.38, 1.5, 0.28, size=18, bold=True, color="blue")
    multiline(s, [
        "让学生拥有明确的阶段目标。",
        "让练习结果产生即时反馈。",
        "让薄弱点修复进入下一步路线。",
        "让教师能够看到过程数据，而非只看到最终分数。"
    ], 7.45, 3.05, 3.8, 1.55, 12.5, "ink")
    footer(s, 5)

    # 6 Mapping
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "教学映射", "游戏元素对应明确教学任务", "爬塔系统的关键在于把游戏机制转化为教学机制。")
    items = [
        ("诊断房", "判断当前知识点掌握情况，决定是否进入训练"),
        ("普通战斗房", "完成基础训练，巩固当前知识点"),
        ("精英战斗房", "处理高难度题目和易错知识点"),
        ("Boss 房", "进行阶段性综合测评"),
        ("补给 / 营火", "错题复盘、知识解释、状态恢复和学习建议"),
        ("奖励系统", "通过成长值、徽章和阶段称号反馈学习成果"),
    ]
    for i, (a, b) in enumerate(items):
        x = 0.85 + (i % 2) * 6.0
        y = 1.85 + (i // 2) * 1.45
        module_box(s, a, b, x, y, 5.25, 1.0, ["blue", "teal", "green", "gold", "red", "blue"][i])
    footer(s, 6)

    # 7 Personalized route
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "个性化路线", "知识图谱、历史答题、掌握度和能力画像共同决定路径", "AI 给出推荐顺序，后端校验合法性，数据库记录全过程。")
    module_box(s, "输入依据", "知识点前置关系、学生正确率、错误类型、耗时、掌握度、能力点得分、节点状态。", 0.85, 2.0, 3.75, 2.1, "blue")
    module_box(s, "AI 作用", "生成推荐路线和薄弱点修复顺序，帮助学生从当前状态进入下一步学习任务。", 4.8, 2.0, 3.75, 2.1, "teal")
    module_box(s, "系统约束", "后端校验课程边界、节点合法性和通关状态；数据库记录 run、node、attempt，保证可追溯。", 8.75, 2.0, 3.75, 2.1, "green")
    card(s, 2.1, 5.15, 9.15, 0.78, fill="blue2", line="blue")
    text(s, "路线设计目标：既体现个性化，又保持业务状态可靠、可解释、可回溯。", 2.45, 5.4, 8.45, 0.2, size=14, bold=True, color="blue", align=PP_ALIGN.CENTER)
    footer(s, 7)

    # 8 AI
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "AI 设计", "AI Agent 嵌入多个业务环节", "AI 能力覆盖内容结构化、学习诊断、评价推荐和教学分析。")
    ai_items = [
        ("内容侧", "知识点抽取、知识关系生成、能力点建议"),
        ("学习侧", "题目提示、错因解释、复习建议"),
        ("评价侧", "主观题、报告类作业和开放答案辅助评价"),
        ("推荐侧", "结合画像和知识图谱生成个性化路径"),
        ("教学侧", "共性错误聚类、风险识别和教学建议"),
    ]
    for i, (a, b) in enumerate(ai_items):
        x = 0.95 + i * 2.45
        module_box(s, a, b, x, 2.05, 2.05, 2.35, ["blue", "teal", "green", "gold", "red"][i])
    card(s, 1.85, 5.35, 9.65, 0.75, fill="paper", line="cyan")
    text(s, "设计原则：AI 给建议，后端做校验，数据库留痕，教师可追溯。", 2.2, 5.6, 8.9, 0.22, size=15, bold=True, color="teal", align=PP_ALIGN.CENTER)
    footer(s, 8)

    # 9 Knowledge
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "图谱设计", "知识图谱与能力图谱", "知识图谱描述知识关系，能力图谱描述课程目标与学生能力成长。")
    module_box(s, "知识图谱", "描述知识点之间的前置、关联、强化和综合应用关系。", 0.95, 2.05, 3.6, 2.2, "blue")
    module_box(s, "能力图谱", "承接课程目标，例如语法理解、文件操作、异常处理、数据分析和项目实践。", 4.85, 2.05, 3.6, 2.2, "teal")
    module_box(s, "映射价值", "学生答题结果可以转化为知识掌握度和能力成长数据。", 8.75, 2.05, 3.6, 2.2, "green")
    card(s, 2.0, 5.15, 9.3, 0.8, fill="cyan2", line="cyan")
    text(s, "系统可回答：错在哪里、对应哪个知识点、影响哪项能力。", 2.35, 5.42, 8.6, 0.22, size=15, bold=True, color="blue", align=PP_ALIGN.CENTER)
    footer(s, 9)

    # 10 Database
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "国产数据库", "KingbaseES V9 与数据模型设计", "数据库围绕课程、学习任务、测评、画像和学情分析建立业务表结构。")
    module_box(s, "业务闭环", "课程、知识点、任务、题目、提交、掌握度、推荐和分析等表共同支撑学习闭环。", 0.9, 2.0, 3.75, 2.25, "blue")
    module_box(s, "国产化适配", "业务 ID 以 VARCHAR 为主，规避国产数据库中整数与字符串比较差异。", 4.8, 2.0, 3.75, 2.25, "teal")
    module_box(s, "扩展设计", "服务层维护逻辑关联，减少外键强绑定，提高扩展性和迁移灵活性。", 8.7, 2.0, 3.75, 2.25, "green")
    card(s, 2.1, 5.2, 9.1, 0.72, fill="paper", line="cyan")
    text(s, "独立 KingbaseES 配置支持国产化部署和测试环境切换。", 2.45, 5.45, 8.4, 0.2, size=14.5, bold=True, color="blue", align=PP_ALIGN.CENTER)
    footer(s, 10)

    # 11 Profile
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "学生画像", "对学习过程进行综合建模", "画像用于学生端学习建议，也用于教师端风险识别和精准干预。")
    profile = [
        ("知识掌握", "知识点掌握度、薄弱点和待复盘节点"),
        ("能力成长", "能力点得分、成长历史和阶段趋势"),
        ("学习表现", "作答结果、错题、耗时和任务完成情况"),
        ("学习投入", "行为记录、连续学习、活跃度和最近活动"),
        ("推荐反馈", "个性化推荐结果、徽章、称号和成就记录"),
    ]
    stack(s, [(a, b, ["blue", "teal", "green", "gold", "red"][i]) for i, (a, b) in enumerate(profile)], 1.35, 1.95, 10.6, 4.55)
    footer(s, 11)

    # 12 Teacher
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "教师端设计", "从结果统计转向过程洞察", "教师端重点服务教学决策，而不是只提供后台管理入口。")
    module_box(s, "班级整体", "客观作答数、错误数、错误率、薄弱知识点和题型表现。", 0.9, 2.0, 3.5, 2.25, "blue")
    module_box(s, "个体差异", "学习状态、学习投入、基础掌握、成长趋势和风险信号。", 4.9, 2.0, 3.5, 2.25, "teal")
    module_box(s, "干预动作", "复习安排、专项练习、课堂讲解和个性化任务分配。", 8.9, 2.0, 3.5, 2.25, "green")
    card(s, 2.0, 5.15, 9.4, 0.8, fill="blue2", line="blue")
    text(s, "教师端定位：教学决策支持中心。", 2.4, 5.43, 8.6, 0.22, size=15, bold=True, color="blue", align=PP_ALIGN.CENTER)
    footer(s, 12)

    # 13 Evaluation
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "测评评价", "测评结果回流到学习闭环", "测评既承担知识掌握判断，也承担学习成果评价和路径调整依据。")
    module_box(s, "客观题", "快速判断知识点掌握情况，更新诊断房与普通战斗房结果。", 0.9, 2.0, 3.5, 2.15, "blue")
    module_box(s, "编程题", "检验语法、流程控制、文件操作和数据处理的综合应用能力。", 4.9, 2.0, 3.5, 2.15, "teal")
    module_box(s, "主观任务", "报告类和开放性答案可结合 AI 进行辅助评价。", 8.9, 2.0, 3.5, 2.15, "green")
    card(s, 1.9, 5.05, 9.55, 0.86, fill="paper", line="cyan")
    text(s, "评价结果进入学生画像和爬塔节点状态，影响后续路线推荐。", 2.25, 5.35, 8.85, 0.22, size=15, bold=True, color="teal", align=PP_ALIGN.CENTER)
    footer(s, 13)

    # 14 Testing
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "质量保障", "自动化测试与国产化测试工具", "测试重点放在爬塔链路、AI 降级链路和教师端分析链路的可靠性。")
    module_box(s, "后端测试", "覆盖业务服务、控制器、AI 服务调用和关键数据处理逻辑。", 0.9, 2.0, 3.5, 2.2, "blue")
    module_box(s, "前端测试", "覆盖 API 封装、路由跳转、游戏资源清单和关键交互逻辑。", 4.9, 2.0, 3.5, 2.2, "teal")
    module_box(s, "国产化工具", "用于接口录制、脚本管理、测试任务、负载测试和性能监控。", 8.9, 2.0, 3.5, 2.2, "green")
    card(s, 2.0, 5.12, 9.4, 0.8, fill="cyan2", line="cyan")
    text(s, "压力测试验证高并发访问、集中答题和教学数据统计场景下的系统稳定性。", 2.35, 5.4, 8.7, 0.22, size=14.2, bold=True, color="blue", align=PP_ALIGN.CENTER)
    footer(s, 14)

    # 15 Promotion position
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "推广应用", "推广应用定位", "后续推广应突出学习激励、个性化学习和学习过程数据分析的特色价值。")
    module_box(s, "模式一：独立智慧课堂平台", "面向课程组、教学团队或实训课程，提供课程建设、资源管理、任务发布、在线测评、学生画像和学情分析等完整能力。适合完整项目展示和小范围课程试点。", 0.95, 2.0, 5.3, 3.25, "blue")
    module_box(s, "模式二：课堂游戏化伴随系统", "在已有课堂平台之外提供学习地图、知识点闯关、成长值、徽章、阶段称号和学习挑战。该模式接入成本更低，也更能体现项目特色。", 7.05, 2.0, 5.3, 3.25, "teal")
    text(s, "推广重点：课堂游戏化伴随系统", 3.1, 6.05, 7.0, 0.26, size=16, bold=True, color="blue", align=PP_ALIGN.CENTER)
    footer(s, 15)

    # 16 Companion
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "伴随系统", "主课堂平台 + 游戏化学习空间", "不改变学校原有教学流程，只为学生端增加更有吸引力的学习入口。")
    module_box(s, "主课堂平台负责", "课程资源、作业发布、测验考试、成绩记录和基础教学管理。", 0.95, 2.0, 5.25, 2.7, "blue")
    module_box(s, "游戏化伴随系统负责", "知识点闯关、学习地图、成长值、徽章称号、阶段挑战、个人成长反馈和班级挑战活动。", 7.1, 2.0, 5.25, 2.7, "teal")
    card(s, 1.85, 5.45, 9.65, 0.85, fill="paper", line="cyan")
    text(s, "完成学习任务 → 获得成长反馈 → 解锁新目标 → 继续学习", 2.25, 5.75, 8.85, 0.22, size=15.5, bold=True, color="blue", align=PP_ALIGN.CENTER)
    footer(s, 16)

    # 17 Promotion value
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "推广价值", "课堂游戏化伴随系统更易落地", "该模式避开整体替换平台的高成本，聚焦学生端体验和学习激励。")
    vals = [
        ("降低门槛", "学校无需放弃原有平台，只需对接账号、课程、任务和学习结果。"),
        ("提升参与", "通过闯关、成长值、徽章和称号，让学生获得持续反馈。"),
        ("多课复用", "课程能拆分为知识点、任务和测评结果，就能生成对应学习地图。"),
        ("差异化强", "避开资源管理和作业发布等同质化功能，突出学生端体验。"),
    ]
    for i, (a, b) in enumerate(vals):
        module_box(s, a, b, 0.95 + (i % 2) * 6.05, 2.0 + (i // 2) * 1.85, 5.25, 1.25, ["blue", "teal", "green", "gold"][i])
    footer(s, 17)

    # 18 Path and effects
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "实施路径", "从项目展示到多平台应用", "推广分阶段推进，先验证完整平台，再抽离游戏化伴随能力。")
    stack(s, [
        ("第一阶段", "完整智慧课堂平台用于课程实训展示、校内课程试点和功能验证", "blue"),
        ("第二阶段", "抽离成长值、徽章、称号、知识点闯关和学习地图，形成伴随系统", "teal"),
        ("第三阶段", "完善统一身份认证、数据同步、任务同步和结果回写，适配多课程多平台", "green"),
    ], 1.0, 1.95, 5.7, 3.55)
    card(s, 7.35, 2.0, 4.65, 3.55, fill="paper", line="cyan")
    text(s, "预期应用效果", 7.75, 2.35, 2.2, 0.28, size=18, bold=True, color="blue")
    multiline(s, [
        "学生：学习过程更加可视化、阶段化和有反馈。",
        "教师：更容易了解进度、薄弱点和参与情况。",
        "平台：补充游戏化学习体验，增强持续使用意愿。",
        "推广：从完整平台转向轻量增强应用，更容易落地。"
    ], 7.75, 3.0, 3.55, 1.55, 12.0, "ink")
    footer(s, 18)

    # 19 Summary
    s = prs.slides.add_slide(blank); bg(s)
    title(s, "项目总结", "项目特色总结", "从学习体验、数据闭环、AI 嵌入和国产化工程四个方面形成项目特色。")
    module_box(s, "学习体验创新", "通过爬塔系统把知识点学习转化为连续挑战。", 0.9, 2.0, 2.8, 2.2, "blue")
    module_box(s, "数据驱动学习", "通过作答、通关、掌握度和行为记录形成学生画像。", 3.95, 2.0, 2.8, 2.2, "teal")
    module_box(s, "AI 深度嵌入", "AI 参与知识组织、学习诊断、个性化推荐和教学建议。", 7.0, 2.0, 2.8, 2.2, "green")
    module_box(s, "国产化落地", "采用 KingbaseES 和国产化测试工具支撑部署与质量保障。", 10.05, 2.0, 2.8, 2.2, "gold")
    card(s, 2.1, 5.3, 9.15, 0.85, fill="cyan2", line="cyan")
    text(s, "让课程知识可组织、学习过程可追踪、能力成长可量化、教学干预可执行。", 2.45, 5.6, 8.45, 0.22, size=15, bold=True, color="blue", align=PP_ALIGN.CENTER)
    footer(s, 19)

    prs.save(OUT_FILE)
    print(OUT_FILE)


if __name__ == "__main__":
    build()
