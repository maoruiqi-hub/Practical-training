from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "output" / "ppt-assets"
OUT = ROOT / "output" / "AI智慧课程平台-学生端教师端功能展示完善版.pptx"


W, H = Inches(13.333), Inches(7.5)
COLORS = {
    "ink": RGBColor(18, 28, 43),
    "muted": RGBColor(95, 111, 130),
    "bg": RGBColor(244, 247, 250),
    "panel": RGBColor(255, 255, 255),
    "teal": RGBColor(20, 127, 145),
    "gold": RGBColor(214, 166, 64),
    "red": RGBColor(202, 73, 61),
    "green": RGBColor(80, 156, 113),
    "blue": RGBColor(76, 137, 220),
    "dark": RGBColor(22, 31, 44),
    "line": RGBColor(221, 228, 236),
    "soft_teal": RGBColor(231, 247, 248),
    "soft_gold": RGBColor(252, 246, 230),
    "soft_blue": RGBColor(235, 243, 255),
    "soft_green": RGBColor(235, 248, 241),
    "soft_red": RGBColor(255, 240, 238),
}


def crop_assets():
    crops = {
        "student_tower_crop.png": ("05-student-tower-pure.png", (0, 0, 1280, 720)),
        "student_battle_crop.png": ("06-student-diagnosis-pure.png", (0, 0, 1280, 720)),
        "teacher_dashboard_crop.png": ("02-teacher-dashboard.png", (0, 0, 1280, 720)),
        "learning_analysis_crop.png": ("03-learning-analysis-clean.png", (200, 70, 1280, 720)),
        "ability_map_crop.png": ("04-ability-map-clean.png", (200, 70, 1280, 720)),
        "login_crop.png": ("01-login.png", (0, 0, 1280, 720)),
    }
    for out_name, (src_name, box) in crops.items():
        src = ASSET_DIR / src_name
        if not src.exists():
            continue
        with Image.open(src) as img:
            img.crop(box).save(ASSET_DIR / out_name)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=18, color="ink", bold=False, align=None, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align or PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color] if isinstance(color, str) else color
    return box


def add_title(slide, title, subtitle=None, dark=False):
    add_text(slide, title, 0.6, 0.45, 8.8, 0.58, size=30, color="panel" if dark else "ink", bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.62, 1.02, 10.8, 0.36, size=12.5, color="muted" if not dark else RGBColor(204, 214, 226))


def add_chip(slide, text, x, y, fill, text_color="panel", w=1.25):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill] if isinstance(fill, str) else fill
    shape.line.color.rgb = shape.fill.fore_color.rgb
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(10.5)
    run.font.bold = True
    run.font.color.rgb = COLORS[text_color] if isinstance(text_color, str) else text_color
    return shape


def add_card(slide, x, y, w, h, title, body, accent="teal", num=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLORS["panel"]
    rect.line.color.rgb = RGBColor(226, 232, 240)
    if num is not None:
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.22), Inches(y + 0.22), Inches(0.42), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLORS[accent]
        badge.line.color.rgb = COLORS[accent]
        add_text(slide, str(num), x + 0.34, y + 0.31, 0.18, 0.16, size=9.5, color="panel", bold=True, align=PP_ALIGN.CENTER)
        tx = x + 0.78
    else:
        tx = x + 0.28
    add_text(slide, title, tx, y + 0.22, w - (tx - x) - 0.25, 0.28, size=14.5, color="ink", bold=True)
    body_h = max(0.28, h - 0.72)
    add_text(slide, body, x + 0.28, y + 0.62, w - 0.56, body_h, size=11.2, color="muted")


def add_band(slide, x, y, w, h, fill="soft_blue"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS[fill]
    return shape


def add_stat(slide, x, y, value, label, fill="soft_teal", accent="teal"):
    add_band(slide, x, y, 2.55, 1.06, fill)
    add_text(slide, value, x + 0.26, y + 0.18, 1.0, 0.36, size=24, color=accent, bold=True)
    add_text(slide, label, x + 0.28, y + 0.66, 1.9, 0.22, size=10.5, color="muted")


def add_section_label(slide, label, x, y, color="teal"):
    add_chip(slide, label, x, y, color, w=1.45)


def add_image(slide, filename, x, y, w, h):
    path = ASSET_DIR / filename
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        add_card(slide, x, y, w, h, "截图占位", f"缺少素材：{filename}", "red")


def add_footer(slide, index):
    add_text(slide, f"{index:02d} / AI智慧课程平台", 11.2, 7.04, 1.55, 0.18, size=8.5, color="muted", align=PP_ALIGN.RIGHT)


def build():
    crop_assets()
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["bg"])
    add_band(slide, 0.62, 0.56, 5.55, 5.95, "panel")
    add_text(slide, "AI智慧课程平台", 0.95, 0.98, 4.7, 0.65, size=38, color="ink", bold=True)
    add_text(slide, "学生端与教师端功能展示", 0.98, 1.78, 4.6, 0.42, size=19, color="teal", bold=True)
    add_text(slide, "面向教育管理与学习辅助场景，平台将课程资源、学习过程、测评反馈、学情分析和 AI Agent 能力整合为一个持续迭代的教学闭环。", 1.0, 2.42, 4.55, 1.12, size=14.5, color="muted")
    add_chip(slide, "课程管理", 1.0, 4.0, "blue", w=1.18)
    add_chip(slide, "知识图谱", 2.35, 4.0, "teal", w=1.18)
    add_chip(slide, "智能批改", 3.7, 4.0, "gold", w=1.18)
    add_chip(slide, "个性化推荐", 1.0, 4.52, "green", w=1.42)
    add_chip(slide, "学情分析", 2.6, 4.52, "red", w=1.18)
    add_image(slide, "login_crop.png", 6.48, 0.72, 5.8, 5.42)
    add_text(slide, "截图基于当前本地前端版本生成，便于后续替换为接入真实数据后的演示画面。", 6.5, 6.34, 5.8, 0.24, size=10.5, color="muted")
    add_footer(slide, 1)

    # 2
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["bg"])
    add_title(slide, "系统架构：三端协同 + 数据闭环", "前端负责场景体验，后端沉淀业务数据，AI Agent 负责知识生成、评价与推荐。")
    add_band(slide, 0.72, 1.55, 11.9, 4.75, "panel")
    layers = [
        ("教师端", "课程建设 / 班级运营 / 学情分析 / 能力图谱 / 学生画像", "blue", 0.95),
        ("学生端", "课程学习 / 任务测验 / 游戏化爬塔 / 错题复习 / 个性推荐", "teal", 0.95),
        ("RESTful API", "课程、资源、题库、任务、提交、进度、画像等统一接口", "gold", 3.0),
        ("AI Agent 服务", "知识图谱构建、智能评分、问题聚类、教学建议、个性化推荐", "green", 4.42),
        ("数据与资源", "关系型数据库、课程资源、学习行为、测评记录、画像数据", "red", 5.42),
    ]
    add_card(slide, 1.0, 2.0, 5.1, 0.95, layers[0][0], layers[0][1], layers[0][2], 1)
    add_card(slide, 7.0, 2.0, 5.1, 0.95, layers[1][0], layers[1][1], layers[1][2], 2)
    for i, item in enumerate(layers[2:], start=3):
        add_card(slide, 2.15, item[3], 8.9, 0.82, item[0], item[1], item[2], i)
    add_text(slide, "闭环逻辑：课程建设产生结构化内容，学生学习产生行为与测评数据，AI 与学情分析生成反馈，再反哺教师教学调整。", 1.12, 6.52, 10.8, 0.35, size=13.5, color="muted")
    add_footer(slide, 2)

    # 3
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["bg"])
    add_title(slide, "双端功能地图：同一套数据服务两类用户", "学生端强调学习路径与反馈，教师端强调组织管理与教学决策。")
    add_section_label(slide, "学生端", 0.78, 1.48, "teal")
    add_section_label(slide, "教师端", 6.92, 1.48, "blue")
    student_features = [
        ("课程学习", "课程列表、课时详情、资源预览、知识图谱导航"),
        ("任务测评", "在线测验、作业提交、系统评分、通关挑战"),
        ("学习反馈", "诊断结果、错题本、学习进度、个人画像"),
        ("个性路径", "基于掌握度、弱点和推荐策略进入下一步学习"),
    ]
    teacher_features = [
        ("课程建设", "课程、课时、资源、知识点、题库统一维护"),
        ("过程管理", "任务发布、提交复核、班级运营、风险处理"),
        ("学情洞察", "错题统计、知识点掌握、题型表现、班级风险"),
        ("智能辅助", "AI 能力图谱、智能评价、问题聚类、教学建议"),
    ]
    for i, (title, body) in enumerate(student_features):
        add_card(slide, 0.78, 2.0 + i * 1.08, 5.35, 0.88, title, body, ["teal", "gold", "green", "blue"][i], i + 1)
    for i, (title, body) in enumerate(teacher_features):
        add_card(slide, 6.92, 2.0 + i * 1.08, 5.35, 0.88, title, body, ["blue", "red", "gold", "green"][i], i + 1)
    add_band(slide, 0.78, 6.42, 11.49, 0.45, "soft_teal")
    add_text(slide, "核心关系：教师配置课程与任务，学生产生学习数据，系统分析数据并生成反馈，教师据此优化后续教学。", 1.0, 6.55, 10.9, 0.18, size=12.3, color="teal", bold=True)
    add_footer(slide, 3)

    # 4
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["bg"])
    add_title(slide, "学生端：游戏化学习入口", "把课程学习路径转化为“爬塔路线”，让知识点学习、诊断、挑战和奖励形成连续体验。")
    add_image(slide, "student_tower_crop.png", 0.68, 1.45, 7.45, 4.85)
    add_stat(slide, 8.45, 1.48, "路径", "知识点节点化", "soft_teal", "teal")
    add_stat(slide, 10.95, 1.48, "状态", "能量/金币/经验", "soft_gold", "gold")
    add_card(slide, 8.45, 2.86, 3.85, 0.98, "路线驱动学习", "章节、知识点和房间节点串成学习路径，学生可以直接进入下一步。", "teal")
    add_card(slide, 8.45, 4.02, 3.85, 0.98, "诊断前置", "进入节点前先确认知识状态，再匹配题包、挑战强度和奖励。", "gold")
    add_card(slide, 8.45, 5.18, 3.85, 0.98, "即时激励", "用角色、血量、奖励和通关结果强化持续学习反馈。", "green")
    add_footer(slide, 4)

    # 5
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["bg"])
    add_title(slide, "学生端：学习、测评与成长闭环", "围绕“学资源、做任务、即时反馈、回看薄弱点”组织学生端能力。")
    add_image(slide, "student_battle_crop.png", 0.68, 1.45, 6.7, 4.95)
    features = [
        ("课程与资源", "查看课程、课时、视频/文档资源，进入知识图谱关联内容。", "blue"),
        ("任务与测验", "完成在线测验、作业提交，系统记录提交状态与得分。", "teal"),
        ("智能诊断", "根据答题结果更新掌握度，提示薄弱知识点和下一步行动。", "gold"),
        ("画像与错题本", "汇总学习进度、能力成长、错题记录，支持个性化复习。", "green"),
    ]
    for i, item in enumerate(features):
        add_card(slide, 7.75, 1.45 + i * 1.25, 4.7, 1.04, item[0], item[1], item[2], i + 1)
    add_footer(slide, 5)

    # 6
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["bg"])
    add_title(slide, "学生端：个性化学习能力", "学生端不只是做题入口，而是围绕学习状态持续给出复习、激励和下一步建议。")
    add_card(slide, 0.78, 1.55, 3.8, 1.12, "学习画像", "汇总学习进度、能力成长、掌握度和学习状态，形成个人学习档案。", "green", 1)
    add_card(slide, 4.82, 1.55, 3.8, 1.12, "错题本", "沉淀错题、薄弱知识点和复习建议，支持针对性补弱。", "red", 2)
    add_card(slide, 8.86, 1.55, 3.4, 1.12, "推荐路径", "根据诊断结果与画像数据，推送下一步学习资源或训练节点。", "teal", 3)
    stages = [
        ("学习行为", "看资源、做任务、答题"),
        ("数据记录", "进度、得分、错题、掌握度"),
        ("智能分析", "薄弱点、风险、推荐"),
        ("个性反馈", "复习路径、挑战节点、成长激励"),
    ]
    for i, (title, body) in enumerate(stages):
        x = 0.92 + i * 3.05
        add_band(slide, x, 3.35, 2.45, 1.6, ["soft_blue", "soft_gold", "soft_teal", "soft_green"][i])
        add_text(slide, f"0{i+1}", x + 0.18, 3.55, 0.55, 0.28, size=16, color=["blue", "gold", "teal", "green"][i], bold=True)
        add_text(slide, title, x + 0.18, 3.92, 1.8, 0.28, size=15, color="ink", bold=True)
        add_text(slide, body, x + 0.18, 4.34, 2.0, 0.34, size=10.8, color="muted")
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.48), Inches(3.95), Inches(0.45), Inches(0.32))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["line"]
            arrow.line.color.rgb = COLORS["line"]
    add_band(slide, 1.1, 5.82, 10.9, 0.68, "panel")
    add_text(slide, "学生端价值：把“课程内容”转化为可执行路径，把“测评结果”转化为可理解反馈，把“薄弱点”转化为下一步行动。", 1.32, 6.05, 10.2, 0.22, size=13.2, color="teal", bold=True)
    add_footer(slide, 6)

    # 7
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["bg"])
    add_title(slide, "教师端：首页聚合待办与风险", "教师进入系统后先看到需要处理的事项，而不是散落在不同菜单里的数据。")
    add_image(slide, "teacher_dashboard_crop.png", 0.68, 1.36, 7.3, 4.72)
    add_card(slide, 8.35, 1.36, 3.95, 1.04, "今日焦点", "汇总待批改、逾期、临近截止等需要优先处理的事项。", "red")
    add_card(slide, 8.35, 2.62, 3.95, 1.04, "待处理队列", "把提交复核、任务异常和课堂运营事项排成教师工作清单。", "teal")
    add_card(slide, 8.35, 3.88, 3.95, 1.04, "异常提醒", "从任务、提交、学情中提取风险信号，辅助教师快速定位问题。", "gold")
    add_card(slide, 8.35, 5.14, 3.95, 1.04, "最近变化", "记录最新提交、临近截止与班级动态，保证教学过程可追踪。", "blue")
    add_footer(slide, 7)

    # 8
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["bg"])
    add_title(slide, "教师端：学情分析", "从测评数据中提取班级共性错题、知识点掌握和教学干预建议。")
    add_image(slide, "learning_analysis_crop.png", 0.68, 1.36, 7.1, 4.72)
    add_stat(slide, 8.18, 1.48, "错题", "共性问题定位", "soft_red", "red")
    add_stat(slide, 10.7, 1.48, "掌握", "知识点状态", "soft_teal", "teal")
    add_card(slide, 8.18, 2.88, 4.5, 0.96, "班级层面", "查看客观题作答数、错误数、班级错误率，快速判断课程阶段性学习质量。", "blue")
    add_card(slide, 8.18, 4.05, 4.5, 0.96, "知识点层面", "按知识点统计错误、作答、掌握度和状态，定位需要复讲或补充训练的内容。", "teal")
    add_card(slide, 8.18, 5.22, 4.5, 0.96, "智能建议", "结合问题聚类和 AI 教学建议，为教师提供可执行的干预方向。", "gold")
    add_footer(slide, 8)

    # 9
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["bg"])
    add_title(slide, "教师端：课程能力图谱与资源组织", "将课程能力点、知识点和学习资源建立映射，支撑后续诊断与推荐。")
    add_image(slide, "ability_map_crop.png", 0.68, 1.36, 7.1, 4.72)
    add_card(slide, 8.18, 1.45, 4.5, 1.0, "能力点维护", "教师可以新增、编辑、删除课程能力点，形成课程能力框架。", "green")
    add_card(slide, 8.18, 2.72, 4.5, 1.0, "知识点映射", "把能力点绑定到知识点与章节，让课程结构可追踪、可解释。", "teal")
    add_card(slide, 8.18, 3.99, 4.5, 1.0, "AI 生成草稿", "通过 AI 生成能力图谱草稿，教师审核后采纳，提升课程建设效率。", "gold")
    add_card(slide, 8.18, 5.26, 4.5, 1.0, "服务推荐", "能力图谱为学生画像、错题分析和个性化学习推荐提供依据。", "blue")
    add_footer(slide, 9)

    # 10
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["bg"])
    add_title(slide, "AI Agent 与教学闭环：从数据到行动", "智能能力贯穿课程建设、测评反馈、学情分析和个性化推荐。")
    ai_items = [
        ("知识图谱构建", "从课程资源中抽取知识点与关系，支撑课程结构化组织。", "teal"),
        ("智能批改", "对作业/测验提交生成评分建议、评价摘要与风险等级。", "gold"),
        ("个性化推荐", "结合画像、错题、掌握度，推荐学习路径与复习重点。", "green"),
        ("教学决策建议", "基于班级错题和薄弱点，生成问题聚类与干预建议。", "blue"),
    ]
    for i, item in enumerate(ai_items):
        x = 0.72 + (i % 2) * 6.1
        y = 1.52 + (i // 2) * 1.62
        add_card(slide, x, y, 5.55, 1.08, item[0], item[1], item[2], i + 1)
    add_band(slide, 0.82, 5.25, 11.6, 1.05, "panel")
    add_text(slide, "平台价值总结", 1.08, 5.5, 2.1, 0.32, size=18, color="ink", bold=True)
    add_text(slide, "教师端获得可执行的班级干预依据，学生端获得更清晰的学习路径与即时反馈，平台侧沉淀可持续迭代的课程知识资产。", 3.25, 5.5, 8.6, 0.42, size=14.2, color="muted")
    add_footer(slide, 10)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
